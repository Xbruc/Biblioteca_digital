import os
from pathlib import Path
import pandas as pd
import streamlit as st
import fiona
import geopandas as gpd
import bcrypt
import streamlit.components.v1 as components
import base64
import hashlib
from pathlib import Path
import streamlit as st
import docx 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import streamlit as st
import mammoth
from docx import Document
import streamlit as st
import mammoth
import io
import pandas as pd
from docx import Document
from pptx import Presentation
from pathlib import Path
import geopandas as gpd
import geopandas as gpd
import fiona
import zipfile
import tempfile


# Habilita suporte a KML
fiona.drvsupport.supported_drivers['KML'] = 'rw'
fiona.drvsupport.supported_drivers['LIBKML'] = 'rw'

#load_dotenv()

#USERS = os.getenv("AUTH_USERS", "").split(",")
#PASSWORD_HASH_ENV = os.getenv("AUTH_PASSWORD_HASH")

USERS = st.secrets["AUTH_USERS"].split(",")
PASSWORD_HASH = st.secrets["PASSWORD_HASH"].encode()

#if not PASSWORD_HASH_ENV:
#    raise RuntimeError("AUTH_PASSWORD_HASH não definida")

#PASSWORD_HASH = PASSWORD_HASH_ENV.encode()

def autenticar():

    if "auth" not in st.session_state:
        st.session_state.auth = False

    if st.session_state.auth:
        return True

    container = st.container()

    with container:
        st.markdown('<div class="login-box">', unsafe_allow_html=True)
        st.markdown('<div class="login-title">🔐 Acesso</div>', unsafe_allow_html=True)

        with st.form("login_form"):

            email = st.text_input("E-mail")
            senha = st.text_input("Senha", type="password")

            submitted = st.form_submit_button("Entrar")

            if submitted:
                if email in USERS and bcrypt.checkpw(
                    senha.encode(),
                    PASSWORD_HASH
                ):
                    st.session_state.auth = True
                    st.rerun()
                else:
                    st.error("Credenciais inválidas")

        st.markdown('</div>', unsafe_allow_html=True)

    st.stop()



# =========================
# CONFIG
# =========================
st.set_page_config(
    page_title="SRH | Biblioteca Digital",
    page_icon="📁",
    layout="wide",
)

st.markdown("""
<style>

/* =====================
   FONTES E BASE
===================== */
html, body, [class*="css"] {
    font-family: "Inter", "Segoe UI", Roboto, sans-serif;
}

/* =====================
   HEADER
===================== */
h1 {
    font-weight: 700;
    letter-spacing: -0.5px;
}

h1 span {
    color: #2563eb;
}

/* =====================
   SIDEBAR
===================== */
section[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #f8fafc 0%, #eef2f7 100%);
    border-right: 1px solid #e5e7eb;
}

section[data-testid="stSidebar"] h2 {
    font-size: 1.1rem;
    font-weight: 600;
}

/* Inputs da sidebar */
section[data-testid="stSidebar"] input,
section[data-testid="stSidebar"] select {
    border-radius: 8px;
}

/* =====================
   EXPANDER (UPLOAD)
===================== */
div[data-testid="stExpander"] {
    border-radius: 10px;
    border: 1px solid #e5e7eb;
    background-color: #ffffff;
}

/* =====================
   LISTA DE ARQUIVOS
===================== */
.file-card {
    background: #ffffff;
    border-radius: 12px;
    padding: 14px 16px;
    margin-bottom: 12px;
    border: 1px solid #e5e7eb;
    box-shadow: 0 1px 3px rgba(0,0,0,0.05);
}

/* Nome do arquivo */
.file-title {
    font-weight: 600;
    font-size: 0.95rem;
}

/* Breadcrumb / diretório */
.file-path {
    font-size: 0.8rem;
    color: #6b7280;
    margin-bottom: 8px;
}

/* =====================
   BOTÕES
===================== */
button {
    border-radius: 8px !important;
    background-color: #2563eb
}

button[kind="primary"] {
    background-color: #2563eb !important;
}

button[kind="secondary"] {
    background-color: #f3f4f6 !important;
}

/* =====================
   ALERTAS
===================== */
.stAlert {
    border-radius: 10px;
}

/* =====================
   DIVIDER
===================== */
hr {
    margin: 2rem 0;
}

/* =====================
   MAPAS / PREVIEW
===================== */
iframe {
    border-radius: 10px;
    border: 1px solid #e5e7eb;
}

</style>
""", unsafe_allow_html=True)

st.caption("Versão v1.0.0 – Ambiente de Testes")
autenticar()


BASE_DIR = Path("data/SRH")
MONTHS = ["JAN","FEV","MAR","ABR","MAI","JUN","JUL","AGO","SET","OUT","NOV","DEZ"]
YEARS  = ["2023", "2024", "2025", "2026", "2027"]

PARECER_TIPOS = [
    "Renovação de outorga",
    "Outorga preventiva",
    "Outorga direito de uso",
]

CARTAS_TIPOS = [
    "Cartas de Pendência",
    "Cartas de Indeferimento",
]

METAS = ["Meta 1", "Meta 2", "Meta 3", "Meta 4", "Meta 5", "Meta 6", "Meta 7"]

AREAS = [
    "Outorgas",
    "Gestão Participativa",
    "Gestão Planejamento",
    "Carta de Inexigibilidade",
    "Manifestação Técnica",
    "Empresas Perfuradoras",
    "Base de Dados",
    "Fiscalização",
    "D Tamponamento",
]

# =========================
# GARANTIR RAÍZES
# =========================
for a in AREAS:
    (BASE_DIR / a.replace(" ", "_")).mkdir(parents=True, exist_ok=True)

# =========================
# INDEXAÇÃO
# =========================
@st.cache_data(show_spinner=False)
def build_index():
    rows = []

    for root, _, files in os.walk(BASE_DIR):
        for f in files:
            full = (Path(root) / f).resolve()  # resolve() evita caminho estranho
            rel  = full.relative_to(BASE_DIR.resolve())
            parts = [p.upper() for p in rel.parts]

            def has(x): return any(x in p for p in parts)

            if has("OUTORGAS"):
                area = "Outorgas"
            elif has("GESTAO_PARTICIPATIVA"):
                area = "Gestão Participativa"
            elif has("GESTAO_PLANEJAMENTO"):
                area = "Gestão Planejamento"
            elif has("CARTA_DE_INEXIGIBILIDADE"):
                area = "Carta de Inexigibilidade"
            elif has("MANIFESTACAO_TECNICA"):
                area = "Manifestação Técnica"
            elif has("EMPRESAS_PERFURADORAS"):
                area = "Empresas Perfuradoras"
            elif has("BASE_DE_DADOS"):
                area = "Base de Dados"
            elif has("FISCALIZACAO"):
                area = "Fiscalização"
            elif has("D_TAMPONAMENTO"):
                area = "D Tamponamento"
            else:
                area = "Outros"

            modalidade = (
                "Outorga Simplificada" if has("OUTORGA_SIMPLIFICADA") else
                "Superficial" if has("SUPERFICIAL") else
                "Subterrânea" if has("SUBTERR") else ""
            )

            ano = next((p for p in parts if p.isdigit() and len(p) == 4), "")
            mes = next((m for m in MONTHS if m in parts), "")

            # Gestão Planejamento: seção e meta
            secao_gp = ""
            meta_gp = ""
            if area == "Gestão Planejamento":
                if has("METAS"):
                    secao_gp = "Metas"
                    meta_gp = next((m for m in METAS if m.upper() in " ".join(parts)), "")
                elif has("RELATORIOS"):
                    secao_gp = "Relatórios"
                elif has("CAPACITACAO"):
                    secao_gp = "Capacitação"


            # Gestão Participativa: eixo e subeixo
            # --- Gestão Participativa (novo modelo por Comitê) ---
            comite_gp = ""
            tipo_doc_gp = ""

            if area == "Gestão Participativa":

                # Comitês
                for p in parts:
                    if p.startswith("BH"):
                        comite_gp = p.replace("BH", "").replace("", " ").title()

                if has("LEI DE CRIACAO"):
                    tipo_doc_gp = "Lei de Criação"
                elif has("OFICIOS"):
                    tipo_doc_gp = "Ofícios"
                elif has("ELEICOES"):
                    tipo_doc_gp = "Eleições"
                elif has("REUNIOES"):
                    tipo_doc_gp = "Reuniões Extra e Ordinárias"
                elif has("TERMOS DE REFERENCIA"):
                    tipo_doc_gp = "Termos de Referência"
                elif has("RELATORIOS"):
                    tipo_doc_gp = "Relatórios"
                elif has("CAPACITACAO"):
                    tipo_doc_gp = "Capacitação"
                elif has("ADMINISTRATIVO"):
                    tipo_doc_gp = "Administrativo"

            rows.append({
                "area": area,
                "comite": comite_gp,
                "tipo_doc": tipo_doc_gp,
                "modalidade": modalidade,
                "tipo": next((t for t in PARECER_TIPOS + CARTAS_TIPOS if t.upper() in " ".join(parts)), ""),
                "ano": ano,
                "mes": mes,
                "secao": secao_gp,
                "meta": meta_gp,
                "arquivo": f,
                "caminho": str(full),
                "relativo": str(rel),
            })

    df = pd.DataFrame(rows)

    # garante que colunas existam mesmo se vazio
    for col in ["area", "comite", "tipo_doc","modalidade","tipo","ano","mes","secao","meta","arquivo","caminho","relativo"]:
        if col not in df.columns:
            df[col] = ""

    return df
df = build_index()
# =========================
# HEADER
# =========================
st.markdown("""
<h1>SRH | Biblioteca Digital</h1>
<p>Consulta hierárquica conforme estrutura institucional.</p>
""", unsafe_allow_html=True)

# =========================
# UPLOAD
# =========================
with st.expander("➕ Adicionar documentos", expanded=False):
    area_up = st.selectbox("Área", AREAS)
    target = None

    if area_up == "Outorgas":
        modalidade_up = st.selectbox("Modalidade", ["Superficial", "Subterrânea", "Outorga Simplificada"])

        if modalidade_up in ["Superficial", "Subterrânea"]:
            documento = st.selectbox("Documento", ["Parecer", "Carta"])
            tipo_up = st.selectbox("Tipo", PARECER_TIPOS if documento == "Parecer" else CARTAS_TIPOS)
            ano_up  = st.selectbox("Ano", YEARS)
            mes_up  = st.selectbox("Mês", MONTHS)
            target = BASE_DIR / "Outorgas" / modalidade_up /documento / tipo_up / ano_up / mes_up

        else:
            categoria_up = st.selectbox("Categoria", ["Arquivo direto", "Cartas de Pendência", "Cartas de Indeferimento"])
            ano_up = st.selectbox("Ano", YEARS)
            mes_up = st.selectbox("Mês", MONTHS)

            if categoria_up == "Arquivo direto":
                target = BASE_DIR / "Outorgas" / "Outorga_Simplificada" / ano_up / mes_up
            else:
                target = BASE_DIR / "Outorgas" / "Outorga_Simplificada" / categoria_up.replace(" ", "_") / ano_up / mes_up

    elif area_up == "Gestão Participativa":

        comite_up = st.selectbox(
            "Comitê de Bacia",
            [
                "BH Munim",
                "BH Mearim",
                "BH Itapecuru",
                "BH Turiacu",
                "BH Pindare",
                "BH Parnaiba",
            ]
        )

        tipo_doc_up = st.selectbox(
            "Tipo de Documento",
            [
                "Lei de Criacao",
                "Oficios",
                "Eleicoes",
                "Reunioes Extra e Ordinarias",
                "Termos de Referencia",
            ]
        )

        target = (
            BASE_DIR
            / "Gestao_Participativa"
            / "Comites de Bacia"
            / comite_up
            / tipo_doc_up)


    elif area_up == "Gestão Planejamento":
        # >>> SEM BACIA <<<
        ano_up = st.selectbox("Ano", YEARS)
        secao_up = st.selectbox("Seção", ["Metas", "Relatórios", "Capacitação"])

        if secao_up == "Metas":
            meta_up = st.selectbox("Meta", METAS)
            target = BASE_DIR / "Gestao_Planejamento" / ano_up / "Metas" / meta_up
        elif secao_up == "Relatórios":
            target = BASE_DIR / "Gestao_Planejamento" / ano_up / "Relatorios"
        else:
            target = BASE_DIR / "Gestao_Planejamento" / ano_up / "Capacitacao"

    elif area_up in ["Carta de Inexigibilidade", "Manifestação Técnica", "Fiscalização", "D Tamponamento"]:
        ano_up = st.selectbox("Ano", YEARS)
        mes_up = st.selectbox("Mês", MONTHS)
        target = BASE_DIR / area_up.replace(" ", "_") / ano_up / mes_up

    else:
        target = BASE_DIR / area_up.replace(" ", "_")

    files = st.file_uploader("Selecionar arquivos", accept_multiple_files=True)

    if files and st.button("Salvar arquivos"):
        target.mkdir(parents=True, exist_ok=True)
        for upf in files:
            with open(target / upf.name, "wb") as out:
                out.write(upf.getbuffer())
        build_index.clear()
        st.success("Arquivos adicionados com sucesso.")
        st.rerun()

# =========================
# SIDEBAR – CONSULTA
# =========================
col1, col2 = st.sidebar.columns([4, 1])
with col1:
    show_dashboard = st.sidebar.button("📊 Dashboard")
with col2:
    if st.sidebar.button("Sair"):
        st.session_state.auth = False
        st.rerun()

st.sidebar.markdown("## 🌐 Busca Global")

global_search = st.sidebar.text_input(
    "Buscar em todo o acervo",
    placeholder="Digite nome do arquivo, comitê, ano..."
)
is_global_search = bool(global_search)




st.sidebar.markdown("---")


area = st.sidebar.selectbox("Área", [""] + AREAS)

ready = False
view = pd.DataFrame()

if area == "Outorgas":
    modalidade = st.sidebar.selectbox("Modalidade", ["", "Superficial", "Subterrânea", "Outorga Simplificada"])

    if modalidade in ["Superficial", "Subterrânea"]:
        documento = st.sidebar.selectbox("Documento", ["", "Parecer", "Carta"])
        tipo = st.sidebar.selectbox("Tipo", [""] + (PARECER_TIPOS if documento == "Parecer" else CARTAS_TIPOS))
        ano  = st.sidebar.selectbox("Ano", [""] + YEARS)
        mes  = st.sidebar.selectbox("Mês", [""] + MONTHS)
        ready = all([modalidade, tipo, ano, mes])
        if ready:
            view = df[
                (df["area"] == "Outorgas") &
                (df["modalidade"] == modalidade) &
                (df["tipo"] == tipo) &
                (df["ano"] == ano) &
                (df["mes"] == mes)
            ]

    elif modalidade == "Outorga Simplificada":
        categoria = st.sidebar.selectbox("Categoria", ["", "Arquivo direto", "Cartas de Pendência", "Cartas de Indeferimento"])
        ano = st.sidebar.selectbox("Ano", [""] + YEARS)
        mes = st.sidebar.selectbox("Mês", [""] + MONTHS)
        ready = all([categoria, ano, mes])
        if ready:
            view = df[
                (df["area"] == "Outorgas") &
                (df["modalidade"] == "Outorga Simplificada") &
                (df["ano"] == ano) &
                (df["mes"] == mes)
            ]
            if categoria == "Arquivo direto":
                view = view[view["tipo"] == ""]
            else:
                view = view[view["tipo"] == categoria]

elif area == "Gestão Participativa":

    # --- Comitês válidos ---
    comites_validos = sorted(
        c for c in df.loc[
            df["area"] == "Gestão Participativa", "comite"
        ].dropna().unique()
        if c.strip()
    )

    comite = st.sidebar.selectbox(
        "Comitê de Bacia",
        [""] + comites_validos
    )

    # --- Tipos documentais válidos (dependem do comitê) ---
    if comite:
        tipos_validos = sorted(
            t for t in df.loc[
                (df["area"] == "Gestão Participativa") &
                (df["comite"] == comite),
                "tipo_doc"
            ].dropna().unique()
            if t.strip()
        )
    else:
        tipos_validos = []

    tipo_doc = st.sidebar.selectbox(
        "Tipo de Documento",
        [""] + tipos_validos
    )

    #  Só fica pronto quando AMBOS forem escolhidos
    ready = bool(comite and tipo_doc)

    if ready:
        view = df[
            (df["area"] == "Gestão Participativa") &
            (df["comite"] == comite) &
            (df["tipo_doc"] == tipo_doc)
        ]



elif area == "Gestão Planejamento":
    # >>> SEM BACIA <<<
    ano = st.sidebar.selectbox("Ano", [""] + YEARS)
    secao = st.sidebar.selectbox("Seção", ["", "Metas", "Relatórios", "Capacitação"])

    if secao == "Metas":
        meta = st.sidebar.selectbox("Meta", [""] + METAS)
        ready = all([ano, secao, meta])
        if ready:
            view = df[
                (df["area"] == "Gestão Planejamento") &
                (df["ano"] == ano) &
                (df["secao"] == "Metas") &
                (df["meta"] == meta)
            ]
    elif secao in ["Relatórios", "Capacitação"]:
        ready = all([ano, secao])
        if ready:
            view = df[
                (df["area"] == "Gestão Planejamento") &
                (df["ano"] == ano) &
                (df["secao"] == secao)
            ]
    else:
        ready = False

elif area:
    # Outras áreas: mantenha a lógica que você já vinha usando
    ready = True
    view = df[df["area"] == area]


# Habilita o driver KML que vem desativado por padrão
fiona.drvsupport.supported_drivers['KML'] = 'rw'

def carregar_kml_kmz(path: Path):
    ext = path.suffix.lower()
    
    if ext == ".kml":
        return gpd.read_file(str(path), driver='KML')
    
    elif ext == ".kmz":
        # KMZ é um ZIP. Precisamos extrair o doc.kml interno
        with zipfile.ZipFile(path, 'r') as zip_ref:
            # Cria uma pasta temporária para extrair
            temp_dir = tempfile.mkdtemp()
            zip_ref.extractall(temp_dir)
            # Busca o arquivo .kml dentro do KMZ
            kml_file = list(Path(temp_dir).glob("*.kml"))[0]
            return gpd.read_file(str(kml_file), driver='KML')

def carregar_kml_kmz(path_obj: Path):
    ext = path_obj.suffix.lower()
    if ext == ".kml":
        return gpd.read_file(str(path_obj), driver='KML')
    elif ext == ".kmz":
        with zipfile.ZipFile(path_obj, 'r') as z:
            temp_dir = tempfile.mkdtemp()
            z.extractall(temp_dir)
            kml_file = next(Path(temp_dir).glob("*.kml"))
            return gpd.read_file(str(kml_file), driver='KML')


def visualizar_arquivo(path: Path, altura=700):
    if not path.exists() or not path.is_file():
        st.error("Arquivo não encontrado.")
        return

    ext = path.suffix.lower()
    file_id = hashlib.md5(str(path).encode()).hexdigest()

    # IMAGENS
    if ext in [".png", ".jpg", ".jpeg", ".gif", ".webp", ".tif", ".tiff"]:
        st.markdown(f"### 🖼 Visualização: {path.name}")
        st.image(str(path), width="stretch") 
        return

    if ext == ".pdf":
        st.markdown(f"### 📄 {path.name}")
    
        try:
            with open(path, "rb") as f:
                pdf_bytes = f.read()
    
            st.pdf(pdf_bytes, height=700)
    
        except Exception as e:
            st.error(f"Erro ao carregar PDF: {e}")
    
        return


    # # WORD
    # if ext in [".docx", ".doc"]:
    #     st.markdown(f"### 📝 Word: {path.name}")
    #     try:
    #         doc = docx.Document(path)
    #         full_text = [para.text for para in doc.paragraphs]
    #         st.markdown('\n\n'.join(full_text))
    #     except Exception as e:
    #         st.error(f"Erro ao ler Word: {e}")
    #     return # Finaliza aqui se for Word
    elif ext == ".docx":
        try:
            # 1. Visualização de ALTA FIDELIDADE (Mammoth)
            with open(path, "rb") as f:
                result = mammoth.convert_to_html(f)
                st.components.v1.html(
                    f"<div style='background:white;padding:15px;color:black;border-radius:10px'>{result.value}</div>", 
                    height=altura, scrolling=True
                )

            # 2. Carregamento do Documento para Edição
            doc = Document(path)
            full_text = "\n".join([p.text for p in doc.paragraphs])
            
            st.subheader("🖋️ Editor de Conteúdo")
            st.info("O texto abaixo será salvo exatamente como editado. Imagens e cabeçalhos originais serão mantidos.")
            texto_editado = st.text_area("Edite, apague ou adicione texto:", value=full_text, height=300)

            if st.button("💾 Validar Edição e Gerar Arquivo"):
                novas_linhas = texto_editado.split('\n')
                total_originais = len(doc.paragraphs)

                for i, linha_texto in enumerate(novas_linhas):
                    if i < total_originais:
                        # --- EDITAR / APAGAR PARÁGRAFO EXISTENTE ---
                        p = doc.paragraphs[i]
                        
                        # Limpamos o texto de todos os runs, mas PROTEGEMOS imagens/objetos
                        texto_ja_inserido = False
                        for run in p.runs:
                            # Verifica se o run é uma imagem ou objeto (equação)
                            is_image = 'drawing' in run._element.xml or 'object' in run._element.xml
                            
                            if not is_image:
                                if not texto_ja_inserido:
                                    run.text = linha_texto # Sobrescreve com o novo texto (mesmo que vazio)
                                    texto_ja_inserido = True
                                else:
                                    run.text = "" # Apaga resíduos de texto antigo nos outros runs
                        
                        # Caso o parágrafo original fosse só imagem e agora tem texto
                        if not texto_ja_inserido and linha_texto.strip() != "":
                            p.add_run(linha_texto)
                    else:
                        # --- ADICIONAR NOVO PARÁGRAFO ---
                        doc.add_paragraph(linha_texto)

                # 3. Processamento do Download
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)

                st.download_button(
                    label="📥 Baixar Word com Alterações",
                    data=buffer,
                    file_name=f"final_{path.name}",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
                st.success("Arquivo preparado com sucesso! O texto antigo foi substituído e o novo foi adicionado.")

        except Exception as e:
            st.error(f"Erro crítico no processamento: {e}")


    # EXCEL (Agora fora do bloco do Word)
    if ext in [".xlsx", ".xls"]:
        st.markdown(f"### 📊 Excel: {path.name}")
        try:
            # Requer: pip install openpyxl
            df_dict = pd.read_excel(path, sheet_name=None, engine='openpyxl')
            aba = st.selectbox("Selecione a aba:", list(df_dict.keys()), key=f"excel_{file_id}")
            st.dataframe(df_dict[aba], width="stretch")
        except Exception as e:
            st.error(f"Erro ao ler Excel: {e}")
        return

    # --- POWERPOINT (.pptx) ---
    if ext == ".pptx":
        st.markdown(f"### 📽️ PowerPoint: {path.name}")
        try:
            prs = Presentation(path)
            for i, slide in enumerate(prs.slides):
                with st.expander(f"Slide {i+1}", expanded=(i == 0)):
                    # Itera por todos os elementos do slide
                    for shape in slide.shapes:
                        # 1. Extrai Texto
                        if hasattr(shape, "text") and shape.text.strip():
                            st.write(shape.text)
                        
                        # 2. Extrai Imagens (Imagens inseridas ou coladas)
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_bytes = shape.image.blob
                            st.image(image_bytes, width="stretch")
                            
                        # 3. Extrai Imagens de Placeholders (Imagens em layouts prontos)
                        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            if hasattr(shape, "image"):
                                st.image(shape.image.blob, width="stretch")
                                
        except Exception as e:
            st.error(f"Erro ao ler PowerPoint: {e}")
        return
    
    #######shape 
    elif ext == ".shp":
        st.markdown(f"### 🗺️ Mapa Shapefile: {path.name}")
        try:
            # 1. Leitura com Geopandas
            gdf = gpd.read_file(str(path))

            # 2. Garantir coordenadas para Mapas Web (WGS84)
            if gdf.crs is not None:
                gdf = gdf.to_crs(epsg=4326)
            
            # 3. Extrair lat/lon para o Streamlit (funciona para Pontos)
            # Se forem Polígonos, pegamos o centro para posicionar o mapa
            gdf['lat'] = gdf.geometry.centroid.y
            gdf['lon'] = gdf.geometry.centroid.x

            # 4. Exibir Tabela de Dados
            with st.expander("📊 Atributos"):
                st.dataframe(gdf.drop(columns='geometry'), width="stretch")

            # 5. Visualização no Mapa do Streamlit
            st.subheader("📍 Localização")
            st.map(gdf[['lat', 'lon']], width="stretch")

        except Exception as e:
            st.error(f"Erro ao ler Shapefile: {e}")
            st.info("Dica: Os arquivos .shx e .dbf precisam estar na mesma pasta.")
        return
    
    elif ext in [".kml", ".kmz"]:
        st.markdown(f"### 🌍 Mapa KML/KMZ: {path.name}")
        try:
            gdf = carregar_kml_kmz(path)

            # 1. Garantir CRS WGS84 para o Streamlit
            if gdf.crs is None:
                gdf.set_crs(epsg=4326, inplace=True)
            else:
                gdf = gdf.to_crs(epsg=4326)

            # 2. Extrair coordenadas para o mapa (Pontos ou Centros de Polígonos)
            gdf['lat'] = gdf.geometry.centroid.y
            gdf['lon'] = gdf.geometry.centroid.x

            # 3. Exibir Tabela e Mapa
            with st.expander("📊 Dados do Arquivo"):
                st.dataframe(gdf.drop(columns='geometry'), use_container_width=True)

            st.subheader("📍 Visualização no Mapa")
            st.map(gdf[['lat', 'lon']], width="stretch")

        except Exception as e:
            st.error(f"Erro ao processar {ext.upper()}: {e}")
            st.info("Dica: Verifique se o arquivo não está corrompido ou se possui camadas incompatíveis.")

        # --- NOVO: KML e KMZ ---
    elif ext in [".kml", ".kmz"]:
        st.markdown(f"### 🌍 Mapa: {path.name}")
        try:
            gdf = carregar_kml_kmz(path)
            # Converte para WGS84 para o Streamlit reconhecer
            if gdf.crs is None: gdf.set_crs(epsg=4326, inplace=True)
            else: gdf = gdf.to_crs(epsg=4326)
            
            # Extrai pontos para o mapa
            gdf['lat'] = gdf.geometry.centroid.y
            gdf['lon'] = gdf.geometry.centroid.x
            
            st.map(gdf[['lat', 'lon']], width="stretch")
            with st.expander("Ver atributos"):
                st.dataframe(gdf.drop(columns='geometry'))
        except Exception as e:
            st.error(f"Erro ao carregar mapa: {e}")
        return


    # OUTROS FORMATOS
    st.info(f"Visualização não disponível para `{ext}`.")



########################### lidando com shp files
import zipfile
import io

def preparar_download_shp(path_shp: Path):
    """Gera um buffer ZIP contendo todos os componentes do Shapefile."""
    buffer = io.BytesIO()
    # Nome base do arquivo (ex: 'mapa' de 'mapa.shp')
    nome_base = path_shp.stem 
    pasta = path_shp.parent

    with zipfile.ZipFile(buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        # Busca todos os arquivos que começam com o mesmo nome na pasta
        for arquivo in pasta.glob(f"{nome_base}.*"):
            zip_file.write(arquivo, arcname=arquivo.name)
    
    buffer.seek(0)
    return buffer

def deletar_shapefile_completo(path_shp: Path):
    """Remove todos os arquivos associados ao Shapefile."""
    nome_base = path_shp.stem
    pasta = path_shp.parent
    
    arquivos_removidos = 0
    for arquivo in pasta.glob(f"{nome_base}.*"):
        try:
            arquivo.unlink() # Deleta o arquivo
            arquivos_removidos += 1
        except Exception as e:
            st.error(f"Erro ao deletar {arquivo.name}: {e}")
    
    return arquivos_removidos

    


def render_lista_arquivos_com_acoes(view_df, build_index_func=None):

    st.session_state.setdefault("preview_path", None)
    st.session_state.setdefault("to_delete_path", None)

    # remove arquivos que já não existem
    view_df = view_df[
        view_df["caminho"].apply(lambda p: Path(p).exists())
    ]

    if view_df.empty:
        st.warning("Nenhum arquivo encontrado.")
        return

    st.markdown(f"### 📂 Arquivos ({len(view_df)})")

    # Callback para definir o caminho no estado imediatamente
    def set_preview_callback(path):
        # Se o caminho clicado for o mesmo que já está em preview, feche-o (None), senão, defina o novo.
        if st.session_state.preview_path == str(path):
            st.session_state.preview_path = None
        else:
            st.session_state.preview_path = str(path)

    for i, r in view_df.reset_index(drop=True).iterrows():

        relativo = r.get("relativo", "")

        if is_global_search and relativo:
            breadcrumb = " › ".join(relativo.split(os.sep)[:-1])
            st.caption(f"📂 {breadcrumb}")


        arquivo = r["arquivo"]
        caminho = Path(r["caminho"])
        relativo = r.get("relativo", "")

        st.markdown(
            f"""
            <div class="file-card">
                <div class="file-title">📄 {arquivo}</div>
                {f'<div class="file-path">📂 {breadcrumb}</div>' if is_global_search and relativo else ''}
            </div>
            """,
            unsafe_allow_html=True
        )

        if relativo:
            st.caption(relativo)

        c1, c2, c3 = st.columns([1.2, 1.0, 1.0])

        # --- LOGICA DE DOWNLOAD (C1) ---
        with c1:
            ext = caminho.suffix.lower()
            
            if ext == ".shp":
                # Chama a função que criamos para zipar o conjunto SHP
                zip_data = preparar_download_shp(caminho) 
                st.download_button(
                    "⬇ Baixar (ZIP)",
                    data=zip_data,
                    file_name=f"{caminho.stem}.zip",
                    key=f"dl_{i}",
                    width="stretch",
                    mime="application/zip"
                )
            else:
                # Download normal para outras extensões
                with open(caminho, "rb") as f:
                    st.download_button(
                        "⬇ Baixar",
                        f,
                        file_name=arquivo,
                        key=f"dl_{i}",
                        width="stretch"
                    )

        # --- LOGICA DE VISUALIZAR (C2) ---
        with c2:
            if caminho.exists():
                button_text = "❌ Fechar Preview" if st.session_state.preview_path == str(caminho) else "👁 Visualizar"
                st.button(
                    button_text,
                    key=f"vw_{i}_{arquivo}", 
                    width="stretch",
                    on_click=set_preview_callback,
                    args=(caminho,) 
                )
            else:
                st.button("—", disabled=True, width="stretch")

        # --- LOGICA DE EXCLUIR (C3) ---
        with c3:
            if st.button("🗑 Apagar", key=f"del_{i}", width="stretch"):

                if ext == ".shp":
                    # Em vez de apenas setar o path, deletamos todos os arquivos do conjunto agora
                    # ou passamos uma flag para o seu gerenciador de exclusão
                    qtd = deletar_shapefile_completo(caminho)
                    st.success(f"Conjunto {caminho.stem} removido!")
                    st.rerun()
                else:
                    # Comportamento padrão para arquivos únicos
                    st.session_state.to_delete_path = str(caminho)

        
        # ================= NOVO: PREVIEW INLINE
        # Se o caminho atual no loop for o mesmo que o selecionado, exiba o visualizador aqui!
        if st.session_state.preview_path == str(caminho):
            # Usamos um container para garantir que a renderização ocorra exatamente neste ponto.
            with st.container():
                visualizar_arquivo(Path(st.session_state.preview_path))


    # Removemos a seção de preview que estava no final, pois agora é inline.
    # st.session_state.preview_path = None # Não limpe o estado aqui!

    # ================= CONFIRMAÇÃO DE EXCLUSÃO
    _confirmar_exclusao_dialog(build_index_func)


def _confirmar_exclusao_dialog(build_index_func=None):

    if not st.session_state.get("to_delete_path"):
        return

    del_path = Path(st.session_state.to_delete_path)

    @st.dialog("Confirmar exclusão")
    def _dlg():
        st.warning("⚠️ Esta ação é permanente.")
        st.write(f"Arquivo: `{del_path.name}`")

        c1, c2 = st.columns(2)

        with c1:
            if st.button("Cancelar", use_container_width=True):
                st.session_state.to_delete_path = None
                st.rerun()

        with c2:
            if st.button("Confirmar exclusão", type="primary", use_container_width=True):
                del_path.unlink(missing_ok=True)

                if build_index_func:
                    build_index_func.clear()

                st.session_state.to_delete_path = None
                st.session_state.preview_path = None
                st.success("Arquivo excluído com sucesso.")
                st.rerun()

    _dlg()



# =========================
# FILTRO GLOBAL (PRIORITÁRIO)
# =========================
if global_search:
    q = global_search.strip().lower()

    view = df[
        df["arquivo"].str.lower().str.contains(q, na=False) |
        df["relativo"].str.lower().str.contains(q, na=False)
    ]

    ready = True

import shutil
from pathlib import Path

def get_directory_size(path: Path):
    total = 0
    for p in path.rglob("*"):
        if p.is_file():
            try:
                total += p.stat().st_size
            except OSError:
                pass
    return total


def render_storage_usage(BASE_DIR):

    st.subheader("💽 Uso de Armazenamento")

    data_dir = BASE_DIR.resolve()

    # 🔹 No Linux (Streamlit Cloud), use "/"
    # 🔹 No Windows local, também funciona
    disk = shutil.disk_usage("/")

    total_disk = disk.total
    free_disk = disk.free
    used_disk = disk.used

    data_size = get_directory_size(data_dir)

    percent_data = (data_size / total_disk) * 100 if total_disk > 0 else 0

    def gb(x):
        return x / (1024 ** 3)

    st.markdown(
        f"""
        **📁 Diretório monitorado:** `{data_dir}`  
        """
    )

    st.progress(min(percent_data / 100, 1.0))

    c1, c2, c3 = st.columns(3)

    with c1:
        st.metric("📦 Data (GB)", f"{gb(data_size):.2f}")

    with c2:
        st.metric("💽 Disco Total (GB)", f"{gb(total_disk):.0f}")

    with c3:
        st.metric("📊 Uso do Data (%)", f"{percent_data:.2f}%")

    st.caption(f"💾 Espaço livre no disco: {gb(free_disk):.1f} GB")

    st.divider()




if show_dashboard:
    ready = False

if show_dashboard:
    render_storage_usage(BASE_DIR)
    st.stop()

# ... depois de montar: ready e view
st.divider()

if not ready:
    st.info("Selecione todos os filtros para listar documentos.")
else:
    # se quiser manter sua mensagem quando vazio, pode deixar:
    if view.empty:
        st.warning("Nenhum arquivo encontrado.")
    else:
         render_lista_arquivos_com_acoes(view)

